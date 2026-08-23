from typing import List
from pptx import Presentation

from app.core.exceptions import DocumentProcessingException
from .base import DocumentLoader, ParsedSection


class PPTXLoader(DocumentLoader):
    """Loads PPTX files, treating each slide as a page."""

    def load(self, file_path: str) -> List[ParsedSection]:
        prs = Presentation(file_path)
        sections: List[ParsedSection] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                sections.append(
                    ParsedSection(
                        content="\n".join(slide_text),
                        page=slide_num
                    )
                )

        if not sections:
            raise DocumentProcessingException("PPTX contained no extractable text.")
            
        return sections